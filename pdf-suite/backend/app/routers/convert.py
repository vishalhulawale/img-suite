"""PDF Conversion endpoints.

Hybrid conversion engine:
  - DOCX/PPTX/XLSX → PDF: LibreOffice headless (high-fidelity)
  - Images → PDF: Pillow (fast, lossless)
  - PDF → Images: PyMuPDF with optimized parallel/streaming rendering
  - PDF → DOCX/XLSX/PPTX: pdf-converter pipeline (high-fidelity layout,
    table detection, image extraction with CPU-only heuristic analysis)
"""

import asyncio
import io
import json
import logging
import os
import struct
import tempfile
import time
import zipfile
import zlib
from concurrent.futures import ThreadPoolExecutor

import fitz  # PyMuPDF
from PIL import Image
from fastapi import APIRouter, BackgroundTasks, File, Form, HTTPException, UploadFile
from fastapi.responses import Response, StreamingResponse

from app.libreoffice import check_libreoffice_installed, convert_with_libreoffice
from app.tasks import ConversionTask, create_task, get_task, load_task_result, remove_task
from app.utils import (
    ALLOWED_IMAGE_MIME,
    ALLOWED_OFFICE_MIME,
    cleanup_dir,
    cleanup_files,
    save_upload,
)

logger = logging.getLogger(__name__)
router = APIRouter()

# Thread pool for parallel page rendering (bounded to avoid OOM)
_render_pool = ThreadPoolExecutor(max_workers=min(4, (os.cpu_count() or 2)))


def _render_page(page_args: tuple) -> tuple:
    """Render a single PDF page to image bytes (runs in thread pool).

    Each call opens/closes its own fitz.Document so there are no
    shared-state conflicts across threads.
    """
    pdf_path, page_index, zoom, output_ext, fmt = page_args
    doc = fitz.open(pdf_path)
    try:
        page = doc[page_index]
        matrix = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        if output_ext == "jpeg":
            # JPEG: use tobytes directly — fast and memory-efficient
            img_data = pix.tobytes(output="jpeg", jpg_quality=90)
        else:
            # PNG: use raw samples + zlib for faster compression than
            # the default libpng path inside MuPDF.
            img_data = _fast_png_encode(pix)
        pix = None  # release pixmap memory early
        return (f"page_{page_index + 1}.{fmt}", img_data)
    finally:
        doc.close()


def _fast_png_encode(pix) -> bytes:
    """Encode a PyMuPDF Pixmap to PNG using fast zlib compression.

    The default `pix.tobytes("png")` uses libpng with compression level 9,
    which is very slow for large images. This uses zlib level 1 (fast)
    which is 3-5x faster with only ~10% larger files.
    """
    w, h = pix.width, pix.height
    samples = pix.samples  # raw RGB bytes

    # Build raw PNG data rows with filter byte 0 (None) per row
    stride = pix.stride
    raw_rows = bytearray()
    for y in range(h):
        raw_rows.append(0)  # filter: None
        raw_rows.extend(samples[y * stride: y * stride + w * pix.n])

    compressed = zlib.compress(bytes(raw_rows), level=1)

    # Assemble PNG file
    def _chunk(chunk_type: bytes, data: bytes) -> bytes:
        c = chunk_type + data
        crc = zlib.crc32(c) & 0xFFFFFFFF
        return struct.pack(">I", len(data)) + c + struct.pack(">I", crc)

    out = b"\x89PNG\r\n\x1a\n"
    # IHDR: width, height, bit_depth=8, color_type=2 (RGB), compression=0, filter=0, interlace=0
    color_type = 2 if pix.n == 3 else 6  # 2=RGB, 6=RGBA
    ihdr_data = struct.pack(">IIBBBBB", w, h, 8, color_type, 0, 0, 0)
    out += _chunk(b"IHDR", ihdr_data)
    out += _chunk(b"IDAT", compressed)
    out += _chunk(b"IEND", b"")
    return out


@router.post("/convert/images")
async def convert_to_images(
    file: UploadFile = File(..., description="PDF file to convert"),
    format: str = Form("png", description="Image format: png or jpg"),
    dpi: int = Form(150, description="Resolution in DPI: 72, 150, 300, or 600"),
):
    """Accept a PDF and kick off an async image conversion task.

    Returns ``{"task_id": "..."}`` immediately.  The client should then
    open an SSE connection on ``/convert/progress/{task_id}`` for real-time
    page-by-page progress, and finally ``GET /convert/download/{task_id}``
    to retrieve the result.
    """
    if format not in ("png", "jpg", "jpeg"):
        raise HTTPException(status_code=400, detail="Format must be png or jpg.")
    if dpi not in (72, 150, 300, 600):
        raise HTTPException(status_code=400, detail="DPI must be 72, 150, 300, or 600.")

    path = await save_upload(file)

    doc = fitz.open(path)
    page_count = doc.page_count
    doc.close()

    task = await create_task(
        total_pages=page_count,
        message=f"Starting conversion: {page_count} page(s) @ {dpi} DPI",
    )
    logger.info(
        "Task %s created – PDF→%s %d page(s) @ %d DPI",
        task.task_id, format.upper(), page_count, dpi,
    )

    # Fire-and-forget background coroutine
    asyncio.create_task(
        _convert_images_worker(task, path, format, dpi, page_count)
    )
    return {"task_id": task.task_id}


async def _convert_images_worker(
    task: ConversionTask,
    path: str,
    format: str,
    dpi: int,
    page_count: int,
) -> None:
    """Background coroutine that renders pages and stores the result."""
    ext = "jpeg" if format in ("jpg", "jpeg") else "png"
    fmt_upper = ext.upper()
    zoom = dpi / 72.0
    t0 = time.perf_counter()
    await task.update(status="processing")

    try:
        loop = asyncio.get_running_loop()
        logger.info(
            "PDF→%s conversion started: %d page(s) @ %d DPI",
            fmt_upper, page_count, dpi,
        )

        if page_count == 1:
            # ── Single page ─────────────────────────────────────────
            page_t0 = time.perf_counter()
            name, data = await loop.run_in_executor(
                _render_pool, _render_page, (path, 0, zoom, ext, format)
            )

            logger.info(
                "  Page 1/1 rendered in %.2fs", time.perf_counter() - page_t0,
            )

            await task.update(
                current_page=1,
                progress=100,
                result_data=data,
                result_filename=name,
                result_media_type="image/png" if format == "png" else "image/jpeg",
            )
        else:
            # ── Multiple pages — parallel thread-pool rendering ─────
            results: list[tuple | None] = [None] * page_count
            completed = 0

            async def _render_one(idx: int) -> None:
                nonlocal completed
                page_t0 = time.perf_counter()
                result = await loop.run_in_executor(
                    _render_pool,
                    _render_page,
                    (path, idx, zoom, ext, format),
                )
                results[idx] = result
                completed += 1
                await task.update(
                    current_page=completed,
                    progress=int(completed / page_count * 100),
                    message=f"Rendered page {completed}/{page_count}",
                )
                logger.info(
                    "  Page %d/%d rendered in %.2fs | elapsed: %.2fs",
                    completed, page_count,
                    time.perf_counter() - page_t0,
                    time.perf_counter() - t0,
                )

            await asyncio.gather(
                *[_render_one(i) for i in range(page_count)]
            )

            zip_buf = io.BytesIO()
            with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_STORED) as zf:
                for name, data in results:  # type: ignore[misc]
                    zf.writestr(name, data)
            zip_buf.seek(0)
            await task.update(
                result_data=zip_buf.getvalue(),
                result_filename="pdf_images.zip",
                result_media_type="application/zip",
            )

        # ── Done ────────────────────────────────────────────────
        elapsed = time.perf_counter() - t0
        size_mb = len(task.result_data) / 1048576 if task.result_data else 0
        await task.update(
            status="complete",
            progress=100,
            completed_at=time.time(),
            message="Conversion complete",
        )
        logger.info(
            "PDF→%s conversion complete: %d page(s) @ %d DPI in %.2fs (%.1f MB)",
            fmt_upper, page_count, dpi, elapsed, size_mb,
        )
    except Exception as e:
        await task.update(
            status="error",
            error=str(e),
            message=f"Conversion failed: {e}",
        )
        logger.exception("Image conversion failed (task %s)", task.task_id)
    finally:
        cleanup_files(path)


# ─── SSE progress stream ────────────────────────────────────────────
@router.get("/convert/progress/{task_id}")
async def conversion_progress(task_id: str):
    """Server-Sent Events stream of real conversion progress."""
    # Verify the task exists before opening the stream
    if not await get_task(task_id):
        raise HTTPException(status_code=404, detail="Task not found")

    async def _event_stream():
        while True:
            task = await get_task(task_id)
            if not task:
                break
            data = {
                "status": task.status,
                "progress": task.progress,
                "currentPage": task.current_page,
                "totalPages": task.total_pages,
                "message": task.message,
            }
            if task.status == "error":
                data["error"] = task.error
            yield f"data: {json.dumps(data)}\n\n"
            if task.status in ("complete", "error"):
                break
            await asyncio.sleep(0.25)

    return StreamingResponse(
        _event_stream(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"},
    )


# ─── Polling status (mobile-friendly fallback for SSE) ──────────────
@router.get("/convert/status/{task_id}")
async def conversion_status(task_id: str):
    """Non-streaming status check for clients that can't use SSE."""
    task = await get_task(task_id)
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    data = {
        "status": task.status,
        "progress": task.progress,
        "currentPage": task.current_page,
        "totalPages": task.total_pages,
        "message": task.message,
    }
    if task.status == "error":
        data["error"] = task.error
    return data


# ─── Download completed result ──────────────────────────────────────
@router.get("/convert/download/{task_id}")
async def download_result(task_id: str):
    """Return the finished conversion result and free server memory."""
    task = await get_task(task_id)
    if not task:
        raise HTTPException(status_code=404, detail="Task not found")
    if task.status == "error":
        raise HTTPException(status_code=500, detail=task.error or "Conversion failed")
    if task.status != "complete":
        raise HTTPException(status_code=409, detail="Task not yet complete")

    await load_task_result(task)
    if not task.result_data:
        raise HTTPException(status_code=410, detail="Result already downloaded")

    data = task.result_data
    media_type = task.result_media_type
    filename = task.result_filename

    # TTL-based expiry in Redis handles cleanup automatically.

    return Response(
        content=data,
        media_type=media_type,
        headers={
            "Content-Disposition": f"attachment; filename={filename}",
            "Content-Length": str(len(data)),
        },
    )


def _is_scanned_pdf(path: str) -> bool:
    """Detect if a PDF is scanned (image-only) by checking text density.

    Returns True if the average characters-per-page is below a threshold
    that indicates the pages are rendered images rather than digital text.
    """
    doc = fitz.open(path)
    try:
        if doc.page_count == 0:
            return False
        total_chars = sum(len(page.get_text("text").strip()) for page in doc)
        avg_chars = total_chars / doc.page_count
        # A typical text page has 200+ characters; scanned pages have near-zero
        return avg_chars < 50
    finally:
        doc.close()


@router.post("/convert/docx")
async def convert_to_docx(
    file: UploadFile = File(..., description="PDF file to convert to DOCX"),
    background_tasks: BackgroundTasks = BackgroundTasks(),
):
    """Convert PDF to DOCX with high-fidelity layout preservation.

    Uses the pdf-converter pipeline for pixel-perfect text, table, and
    image extraction with intelligent layout analysis.
    """
    from app.pdf_converter.config import OutputFormat, PipelineConfig
    from app.pdf_converter.pipeline import ConversionPipeline

    path = await save_upload(file)
    background_tasks.add_task(cleanup_files, path)

    # Scanned PDF detection — fail fast with clear guidance
    if _is_scanned_pdf(path):
        raise HTTPException(
            status_code=422,
            detail=(
                "This PDF appears to be scanned (image-only). "
                "Text extraction requires OCR processing. "
                "Please use an OCR tool first, then retry conversion."
            ),
        )

    tmp_docx = None
    try:
        config = PipelineConfig(output_format=OutputFormat.DOCX)
        config.apply_cpu_profile()
        pipeline = ConversionPipeline(config)

        tmp_docx = tempfile.mktemp(suffix=".docx", prefix="convert_")

        loop = asyncio.get_running_loop()
        output_path = await loop.run_in_executor(
            None, pipeline.convert, path, tmp_docx,
        )

        with open(output_path, "rb") as f:
            data = f.read()

        return Response(
            content=data,
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={
                "Content-Disposition": "attachment; filename=converted.docx",
                "Content-Length": str(len(data)),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("DOCX conversion failed")
        raise HTTPException(status_code=500, detail=f"DOCX conversion failed: {e}")
    finally:
        if tmp_docx and os.path.exists(tmp_docx):
            try:
                os.unlink(tmp_docx)
            except OSError:
                pass


@router.post("/convert/xlsx")
async def convert_to_xlsx(
    file: UploadFile = File(..., description="PDF file to convert to XLSX"),
    background_tasks: BackgroundTasks = BackgroundTasks(),
):
    """Convert PDF to XLSX with intelligent table extraction and formatting."""
    from app.pdf_converter.config import OutputFormat, PipelineConfig
    from app.pdf_converter.pipeline import ConversionPipeline

    path = await save_upload(file)
    background_tasks.add_task(cleanup_files, path)

    tmp_xlsx = None
    try:
        config = PipelineConfig(output_format=OutputFormat.XLSX)
        config.apply_cpu_profile()
        pipeline = ConversionPipeline(config)

        tmp_xlsx = tempfile.mktemp(suffix=".xlsx", prefix="convert_")

        loop = asyncio.get_running_loop()
        output_path = await loop.run_in_executor(
            None, pipeline.convert, path, tmp_xlsx,
        )

        with open(output_path, "rb") as f:
            data = f.read()

        return Response(
            content=data,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={
                "Content-Disposition": "attachment; filename=converted.xlsx",
                "Content-Length": str(len(data)),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("XLSX conversion failed")
        raise HTTPException(status_code=500, detail=f"XLSX conversion failed: {e}")
    finally:
        if tmp_xlsx and os.path.exists(tmp_xlsx):
            try:
                os.unlink(tmp_xlsx)
            except OSError:
                pass


@router.post("/convert/pptx")
async def convert_to_pptx(
    file: UploadFile = File(..., description="PDF file to convert to PPTX"),
    background_tasks: BackgroundTasks = BackgroundTasks(),
):
    """Convert PDF to PPTX with pixel-perfect text, table, and image positioning."""
    from app.pdf_converter.config import OutputFormat, PipelineConfig
    from app.pdf_converter.pipeline import ConversionPipeline

    path = await save_upload(file)
    background_tasks.add_task(cleanup_files, path)

    tmp_pptx = None
    try:
        config = PipelineConfig(output_format=OutputFormat.PPTX)
        config.apply_cpu_profile()
        pipeline = ConversionPipeline(config)

        tmp_pptx = tempfile.mktemp(suffix=".pptx", prefix="convert_")

        loop = asyncio.get_running_loop()
        output_path = await loop.run_in_executor(
            None, pipeline.convert, path, tmp_pptx,
        )

        with open(output_path, "rb") as f:
            data = f.read()

        return Response(
            content=data,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={
                "Content-Disposition": "attachment; filename=converted.pptx",
                "Content-Length": str(len(data)),
            },
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("PPTX conversion failed")
        raise HTTPException(status_code=500, detail=f"PPTX conversion failed: {e}")
    finally:
        if tmp_pptx and os.path.exists(tmp_pptx):
            try:
                os.unlink(tmp_pptx)
            except OSError:
                pass


@router.post("/convert/to-pdf")
async def convert_to_pdf(
    file: UploadFile = File(..., description="File to convert to PDF (DOCX, PPTX, XLSX, or Image)"),
    background_tasks: BackgroundTasks = BackgroundTasks(),
):
    """Convert DOCX, PPTX, XLSX, or images to PDF.

    Routing:
      - .docx / .pptx / .xlsx → LibreOffice headless (high-fidelity)
      - .png / .jpg / .jpeg / .webp → Pillow
    Falls back to basic text extraction if LibreOffice is unavailable.
    """
    allowed = ALLOWED_IMAGE_MIME | ALLOWED_OFFICE_MIME
    path = await save_upload(file, allowed_mimes=allowed)
    background_tasks.add_task(cleanup_files, path)

    ext = os.path.splitext(file.filename or "")[1].lower()

    try:
        if ext in (".png", ".jpg", ".jpeg", ".webp"):
            return _image_to_pdf(path)
        elif ext in (".docx", ".pptx", ".xlsx"):
            return await _office_to_pdf(path, background_tasks)
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")
    except HTTPException:
        raise
    except Exception as e:
        logger.exception("Conversion to PDF failed")
        raise HTTPException(status_code=500, detail=f"Conversion to PDF failed: {e}")


# ─── Helper functions for → PDF conversion ────────────────


def _image_to_pdf(path: str) -> StreamingResponse:
    """Convert an image to PDF using Pillow."""
    img = Image.open(path)
    if img.mode in ("RGBA", "P"):
        img = img.convert("RGB")
    buf = io.BytesIO()
    img.save(buf, "PDF", resolution=150)
    data = buf.getvalue()
    return Response(
        content=data,
        media_type="application/pdf",
        headers={
            "Content-Disposition": "attachment; filename=converted.pdf",
            "Content-Length": str(len(data)),
        },
    )


async def _office_to_pdf(path: str, background_tasks: BackgroundTasks) -> StreamingResponse:
    """Convert DOCX/PPTX/XLSX → PDF via LibreOffice headless.

    Falls back to basic PyMuPDF text extraction if LibreOffice is not
    installed, with a warning header so callers know quality is degraded.
    """
    if not check_libreoffice_installed():
        logger.warning(
            "LibreOffice not found — falling back to basic conversion. "
            "Install LibreOffice for high-fidelity output."
        )
        return _basic_text_to_pdf(path)

    output_path = await convert_with_libreoffice(path, output_format="pdf")
    work_dir = os.path.dirname(output_path)
    background_tasks.add_task(cleanup_dir, work_dir)

    with open(output_path, "rb") as f:
        data = f.read()

    return Response(
        content=data,
        media_type="application/pdf",
        headers={
            "Content-Disposition": "attachment; filename=converted.pdf",
            "Content-Length": str(len(data)),
        },
    )


def _basic_text_to_pdf(path: str) -> StreamingResponse:
    """Minimal fallback: extract raw text and render to PDF.

    Only used when LibreOffice is unavailable. Produces low-fidelity output
    and sets an X-Conversion-Quality header to signal the degradation.
    """
    ext = os.path.splitext(path)[1].lower()
    pdf_doc = fitz.open()

    lines: list[str] = []

    if ext == ".docx":
        from docx import Document as DocxDocument
        docx_doc = DocxDocument(path)
        lines = [p.text for p in docx_doc.paragraphs if p.text.strip()]
    elif ext == ".pptx":
        from pptx import Presentation as PptxPresentation
        prs = PptxPresentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            lines.append(para.text.strip())
            lines.append("")  # slide separator
    elif ext == ".xlsx":
        from openpyxl import load_workbook
        wb = load_workbook(path, data_only=True)
        for ws in wb.worksheets:
            lines.append(f"--- {ws.title} ---")
            for row in ws.iter_rows(min_row=1, values_only=True):
                text = "  |  ".join(str(c) if c is not None else "" for c in row)
                if text.strip().replace("|", "").strip():
                    lines.append(text)

    page = pdf_doc.new_page(width=595, height=842)
    y = 50
    for line in lines:
        length = fitz.get_text_length(line, fontname="helv", fontsize=11)
        num_lines = max(1, int(length / 495) + 1)
        block_h = num_lines * 16

        if y + block_h > 792:
            page = pdf_doc.new_page(width=595, height=842)
            y = 50

        page.insert_textbox(
            fitz.Rect(50, y, 545, y + block_h + 10),
            line, fontname="helv", fontsize=11, color=(0, 0, 0),
        )
        y += block_h + 4

    buf = io.BytesIO()
    pdf_doc.save(buf)
    pdf_doc.close()
    data = buf.getvalue()

    return Response(
        content=data,
        media_type="application/pdf",
        headers={
            "Content-Disposition": "attachment; filename=converted.pdf",
            "Content-Length": str(len(data)),
            "X-Conversion-Quality": "low-fidelity",
        },
    )
