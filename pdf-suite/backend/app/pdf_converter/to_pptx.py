"""PPTX converter — high-fidelity PDF to PowerPoint conversion.

Each PDF page becomes a slide with:
  - Exact text box positioning and formatting
  - Images at original resolution and position
  - Tables with formatting
  - Slide dimensions matching PDF page size
"""

from __future__ import annotations

import io
import logging

from pptx import Presentation
from pptx.util import Pt, Inches, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

from app.pdf_converter.config import PipelineConfig
from app.pdf_converter.models import (
    BorderStyle, Color, Document, DrawingElement, ElementType, FontInfo, ImageElement,
    Page, PageElement, PageType, Rect, TableCell, TableElement, TextAlignment,
    TextBlock, TextLine, TextSpan,
)

logger = logging.getLogger(__name__)

PT_TO_EMU = 12700


def convert_to_pptx(doc: Document, output_path: str, config: PipelineConfig) -> str:
    """Convert document model to PPTX."""
    prs = Presentation()

    if doc.pages:
        ref_page = doc.pages[0]
        prs.slide_width = Emu(int(ref_page.width * PT_TO_EMU))
        prs.slide_height = Emu(int(ref_page.height * PT_TO_EMU))

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for page in doc.pages:
        slide = prs.slides.add_slide(blank_layout)

        if page.background_color:
            bg = page.background_color
            if bg.r < 0.99 or bg.g < 0.99 or bg.b < 0.99:
                _set_slide_background(slide, bg)

        elements = sorted(page.elements, key=lambda e: (e.z_order, e.reading_order))

        image_bboxes = [
            e.bbox for e in elements
            if e.element_type == ElementType.IMAGE and not e.is_background
        ]

        bg_draw_bboxes = [
            e for e in elements
            if e.element_type == ElementType.DRAWING and e.is_background
            and isinstance(e.element, DrawingElement) and e.element.fill_color is not None
        ]

        for elem in elements:
            if elem.element_type == ElementType.IMAGE and elem.is_background:
                _add_image(slide, elem, page, config)

        for elem in elements:
            if elem.element_type == ElementType.DRAWING and elem.is_background:
                _add_drawing(slide, elem, page, config)

        for elem in elements:
            if elem.element_type == ElementType.IMAGE and not elem.is_background:
                if _is_occluded_by_bg_drawing(elem, bg_draw_bboxes):
                    continue
                _add_image(slide, elem, page, config)

        for elem in elements:
            if elem.element_type == ElementType.DRAWING and not elem.is_background:
                if _is_artifact_drawing(elem, image_bboxes):
                    continue
                _add_drawing(slide, elem, page, config)

        for elem in elements:
            if elem.element_type == ElementType.TEXT_BLOCK and not elem.is_background:
                if _is_occluded_by_bg_drawing(elem, bg_draw_bboxes):
                    continue
                _add_text_block(slide, elem, page, config)
            elif elem.element_type == ElementType.TABLE and not elem.is_background:
                _add_table(slide, elem, page, config)

    prs.save(output_path)
    logger.info("PPTX saved to: %s", output_path)
    return output_path


# ---------------------------------------------------------------------------
# Slide background
# ---------------------------------------------------------------------------

def _set_slide_background(slide, color: Color):
    """Set solid color slide background."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*color.to_rgb_int())


# ---------------------------------------------------------------------------
# Text blocks
# ---------------------------------------------------------------------------

def _add_text_block(
    slide,
    elem: PageElement,
    page: Page,
    config: PipelineConfig,
):
    """Add a text block as a positioned text box on the slide."""
    block = elem.element
    if not isinstance(block, TextBlock):
        return

    if not block.text.strip():
        return

    left = Emu(int(elem.bbox.x0 * PT_TO_EMU))
    top = Emu(int(elem.bbox.y0 * PT_TO_EMU))
    raw_width = max(elem.bbox.width, 10)
    max_width = max(page.width - elem.bbox.x0, raw_width)
    max_font_size = max(
        (s.font.size for l in block.lines for s in l.spans),
        default=12.0,
    )
    width_factor = 1.25 if max_font_size >= 18 else 1.15
    width = Emu(int(min(raw_width * width_factor, max_width) * PT_TO_EMU))
    height = Emu(int(max(elem.bbox.height, 10) * PT_TO_EMU))

    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE

    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)

    if block.background_color:
        fill = txbox.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*block.background_color.to_rgb_int())

    first_para = True

    for line_idx, line in enumerate(block.lines):
        if first_para:
            para = tf.paragraphs[0]
            first_para = False
        else:
            para = tf.add_paragraph()

        align_map = {
            TextAlignment.LEFT: PP_ALIGN.LEFT,
            TextAlignment.CENTER: PP_ALIGN.CENTER,
            TextAlignment.RIGHT: PP_ALIGN.RIGHT,
            TextAlignment.JUSTIFY: PP_ALIGN.JUSTIFY,
        }
        para.alignment = align_map.get(block.alignment, PP_ALIGN.LEFT)

        if block.line_spacing > 0 and block.line_spacing != 1.0:
            para.line_spacing = Pt(block.line_spacing * (line.dominant_font.size if line.dominant_font else 12))

        if block.space_before > 0:
            para.space_before = Pt(block.space_before)
        if block.space_after > 0:
            para.space_after = Pt(block.space_after)

        for span in line.spans:
            _add_span_to_para(para, span)


def _add_span_to_para(para, span: TextSpan):
    """Add a TextSpan as a run to a PPTX paragraph."""
    run = para.add_run()
    run.text = span.text

    font = run.font

    font.name = span.font.name
    font.size = Pt(span.font.size)
    font.bold = span.font.bold
    font.italic = span.font.italic
    font.underline = span.font.underline

    if span.font.strikethrough:
        from pptx.oxml.ns import qn
        font._element.set("strike", "sngStrike")

    if span.font.color:
        font.color.rgb = RGBColor(*span.font.color.to_rgb_int())


# ---------------------------------------------------------------------------
# Images
# ---------------------------------------------------------------------------

def _add_image(
    slide,
    elem: PageElement,
    page: Page,
    config: PipelineConfig,
):
    """Add an image at its exact PDF position."""
    image = elem.element
    if not isinstance(image, ImageElement):
        return

    try:
        img_stream = io.BytesIO(image.image_data)

        left = Emu(int(elem.bbox.x0 * PT_TO_EMU))
        top = Emu(int(elem.bbox.y0 * PT_TO_EMU))
        width = Emu(int(elem.bbox.width * PT_TO_EMU))
        height = Emu(int(elem.bbox.height * PT_TO_EMU))

        slide.shapes.add_picture(img_stream, left, top, width, height)

    except Exception as e:
        logger.warning("Failed to add image to slide: %s", e)


# ---------------------------------------------------------------------------
# Tables
# ---------------------------------------------------------------------------

def _add_table(
    slide,
    elem: PageElement,
    page: Page,
    config: PipelineConfig,
):
    """Add a table shape at its exact PDF position."""
    table = elem.element
    if not isinstance(table, TableElement):
        return

    if table.num_rows == 0 or table.num_cols == 0:
        return

    left = Emu(int(elem.bbox.x0 * PT_TO_EMU))
    top = Emu(int(elem.bbox.y0 * PT_TO_EMU))
    width = Emu(int(elem.bbox.width * PT_TO_EMU))
    height = Emu(int(elem.bbox.height * PT_TO_EMU))

    shape = slide.shapes.add_table(
        table.num_rows, table.num_cols, left, top, width, height
    )
    pptx_table = shape.table

    for col_idx, col_width in enumerate(table.col_widths):
        if col_idx < len(pptx_table.columns):
            pptx_table.columns[col_idx].width = Emu(int(col_width * PT_TO_EMU))

    for row_idx, row_height in enumerate(table.row_heights):
        if row_idx < len(pptx_table.rows):
            pptx_table.rows[row_idx].height = Emu(int(row_height * PT_TO_EMU))

    merged: set[tuple[int, int]] = set()

    for cell in table.cells:
        if (cell.row_index, cell.col_index) in merged:
            continue
        if cell.row_index >= table.num_rows or cell.col_index >= table.num_cols:
            continue

        pptx_cell = pptx_table.cell(cell.row_index, cell.col_index)

        if cell.row_span > 1 or cell.col_span > 1:
            end_row = min(cell.row_index + cell.row_span - 1, table.num_rows - 1)
            end_col = min(cell.col_index + cell.col_span - 1, table.num_cols - 1)
            try:
                other = pptx_table.cell(end_row, end_col)
                pptx_cell.merge(other)
            except Exception:
                pass

            for r in range(cell.row_index, end_row + 1):
                for c in range(cell.col_index, end_col + 1):
                    if (r, c) != (cell.row_index, cell.col_index):
                        merged.add((r, c))

        _populate_pptx_cell(pptx_cell, cell)

        if cell.background_color:
            fill = pptx_cell.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(*cell.background_color.to_rgb_int())

        va_map = {
            "top": MSO_ANCHOR.TOP,
            "center": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }
        pptx_cell.vertical_anchor = va_map.get(cell.vertical_alignment, MSO_ANCHOR.TOP)

        pptx_cell.margin_left = Emu(int(cell.padding_left * PT_TO_EMU))
        pptx_cell.margin_right = Emu(int(cell.padding_right * PT_TO_EMU))
        pptx_cell.margin_top = Emu(int(cell.padding_top * PT_TO_EMU))
        pptx_cell.margin_bottom = Emu(int(cell.padding_bottom * PT_TO_EMU))


def _populate_pptx_cell(pptx_cell, cell: TableCell):
    """Fill cell content with formatted text."""
    tf = pptx_cell.text_frame
    tf.word_wrap = True

    if tf.paragraphs:
        tf.paragraphs[0].clear()

    first_para = True
    for block in cell.content:
        for line_idx, line in enumerate(block.lines):
            if first_para:
                para = tf.paragraphs[0]
                first_para = False
            else:
                para = tf.add_paragraph()

            for span in line.spans:
                _add_span_to_para(para, span)


# ---------------------------------------------------------------------------
# Drawings (vector shapes)
# ---------------------------------------------------------------------------

def _add_drawing(
    slide,
    elem: PageElement,
    page: Page,
    config: PipelineConfig,
):
    """Add a drawing element as a positioned rectangle shape."""
    from pptx.enum.shapes import MSO_SHAPE

    drawing = elem.element
    if not isinstance(drawing, DrawingElement):
        return
    if not drawing.fill_color and not drawing.stroke_color:
        return

    left = Emu(int(elem.bbox.x0 * PT_TO_EMU))
    top = Emu(int(elem.bbox.y0 * PT_TO_EMU))
    width = Emu(int(max(elem.bbox.width, 1) * PT_TO_EMU))
    height = Emu(int(max(elem.bbox.height, 1) * PT_TO_EMU))

    try:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)

        from pptx.oxml.ns import qn as pptx_qn
        style_elem = shape._element.find(pptx_qn('p:style'))
        if style_elem is not None:
            shape._element.remove(style_elem)

        if drawing.fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*drawing.fill_color.to_rgb_int())
            if drawing.opacity < 0.999:
                from pptx.oxml.ns import qn as pptx_qn
                fill_elem = shape.fill._fill
                srgbClr = fill_elem.find('.//' + pptx_qn('a:srgbClr'))
                if srgbClr is not None:
                    from pptx.oxml import OxmlElement as PptxOxmlElement
                    alpha = PptxOxmlElement('a:alpha')
                    alpha.set('val', str(int(drawing.opacity * 100000)))
                    srgbClr.append(alpha)
        else:
            shape.fill.background()

        if drawing.stroke_color:
            shape.line.color.rgb = RGBColor(*drawing.stroke_color.to_rgb_int())
            shape.line.width = Pt(drawing.stroke_width)
        else:
            shape.line.width = Emu(0)

        shape.text_frame.word_wrap = False
    except Exception as e:
        logger.warning("Failed to add drawing shape: %s", e)


def _is_artifact_drawing(elem: PageElement, image_bboxes: list) -> bool:
    """Check if a drawing is a small UI artifact fully contained within an image."""
    if elem.bbox.area <= 0:
        return True
    if elem.bbox.width > 200 or elem.bbox.height > 100:
        return False
    for img_bbox in image_bboxes:
        if img_bbox.contains(elem.bbox):
            return True
    return False


def _is_occluded_by_bg_drawing(elem: PageElement, bg_draw_elems: list[PageElement]) -> bool:
    """Check if a FG element is hidden behind an opaque background drawing."""
    if not bg_draw_elems:
        return False
    eb = elem.bbox
    if eb.area <= 0:
        return False

    for bg_elem in bg_draw_elems:
        draw = bg_elem.element
        fc = draw.fill_color
        if fc is None:
            continue
        if fc.r > 0.9 and fc.g > 0.9 and fc.b > 0.9:
            continue

        db = bg_elem.bbox
        inter = eb.intersection(db)
        if inter is None:
            continue
        overlap = inter.area / eb.area if eb.area > 0 else 0
        if overlap < 0.7:
            continue

        if elem.element_type == ElementType.TEXT_BLOCK and isinstance(elem.element, TextBlock):
            block = elem.element
            colors = [s.font.color for l in block.lines for s in l.spans if s.font.color]
            if colors:
                avg_brightness = sum(c.r + c.g + c.b for c in colors) / (3 * len(colors))
                if avg_brightness > 0.7:
                    continue
            return True

        if elem.element_type == ElementType.IMAGE:
            return True

    return False
